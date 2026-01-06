"""
BRO File Conversion Tools
Convert between file formats - images, PDFs, documents, presentations.
All processing happens locally with no cloud uploads.
"""

import os
import sys
from typing import Optional, List
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from .registry import tool

# Import libraries with graceful fallbacks
try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


# ============================================================================
# IMAGE CONVERSION TOOLS
# ============================================================================

@tool("convert_image", "Converts an image to another format. Supports PNG, JPG, BMP, GIF, WEBP, ICO.")
def convert_image(input_path: str, output_format: str, output_path: str = None) -> str:
    """
    Convert an image to a different format.
    
    Args:
        input_path: Path to the source image
        output_format: Target format (png, jpg, bmp, gif, webp, ico)
        output_path: Optional output path (default: same name, new extension)
        
    Returns:
        Path to the converted image
    """
    if not PILLOW_AVAILABLE:
        return "❌ Pillow not installed. Run: pip install Pillow"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    output_format = output_format.lower().strip('.')
    format_map = {
        'jpg': 'JPEG', 'jpeg': 'JPEG',
        'png': 'PNG',
        'bmp': 'BMP',
        'gif': 'GIF',
        'webp': 'WEBP',
        'ico': 'ICO',
        'tiff': 'TIFF', 'tif': 'TIFF'
    }
    
    if output_format not in format_map:
        return f"❌ Unsupported format: {output_format}. Supported: {', '.join(format_map.keys())}"
    
    try:
        image = Image.open(input_path)
        
        # Handle transparency for formats that don't support it
        if output_format in ['jpg', 'jpeg', 'bmp'] and image.mode in ['RGBA', 'LA', 'P']:
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = background
        
        # Generate output path if not provided
        if not output_path:
            base_name = os.path.splitext(input_path)[0]
            output_path = f"{base_name}.{output_format}"
        
        # Save with appropriate settings
        save_kwargs = {}
        if output_format in ['jpg', 'jpeg']:
            save_kwargs['quality'] = 90
        elif output_format == 'webp':
            save_kwargs['quality'] = 85
        
        image.save(output_path, format_map[output_format], **save_kwargs)
        return f"✓ Converted: {output_path}"
    
    except Exception as e:
        return f"❌ Conversion failed: {str(e)}"


@tool("resize_image", "Resizes an image to specified dimensions.")
def resize_image(input_path: str, width: int = None, height: int = None, output_path: str = None) -> str:
    """
    Resize an image, maintaining aspect ratio if only one dimension specified.
    
    Args:
        input_path: Path to the source image
        width: Target width (optional)
        height: Target height (optional)
        output_path: Optional output path (default: overwrites original)
        
    Returns:
        Path to the resized image
    """
    if not PILLOW_AVAILABLE:
        return "❌ Pillow not installed"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    if width is None and height is None:
        return "❌ Specify at least width or height"
    
    try:
        image = Image.open(input_path)
        orig_width, orig_height = image.size
        
        # Calculate dimensions maintaining aspect ratio
        if width and not height:
            ratio = width / orig_width
            height = int(orig_height * ratio)
        elif height and not width:
            ratio = height / orig_height
            width = int(orig_width * ratio)
        
        resized = image.resize((width, height), Image.Resampling.LANCZOS if hasattr(Image, 'Resampling') else Image.LANCZOS)
        
        save_path = output_path or input_path
        resized.save(save_path)
        
        return f"✓ Resized to {width}x{height}: {save_path}"
    
    except Exception as e:
        return f"❌ Resize failed: {str(e)}"


@tool("compress_image", "Compresses an image to reduce file size.")
def compress_image(input_path: str, quality: int = 70, output_path: str = None) -> str:
    """
    Compress an image to reduce file size.
    
    Args:
        input_path: Path to the source image
        quality: Compression quality 1-100 (lower = smaller file)
        output_path: Optional output path
        
    Returns:
        Path to compressed image with size info
    """
    if not PILLOW_AVAILABLE:
        return "❌ Pillow not installed"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    try:
        original_size = os.path.getsize(input_path)
        image = Image.open(input_path)
        
        # Convert to RGB if needed
        if image.mode in ['RGBA', 'P']:
            image = image.convert('RGB')
        
        save_path = output_path or input_path
        image.save(save_path, 'JPEG', quality=quality, optimize=True)
        
        new_size = os.path.getsize(save_path)
        reduction = ((original_size - new_size) / original_size) * 100
        
        return f"✓ Compressed: {save_path} ({original_size//1024}KB → {new_size//1024}KB, {reduction:.1f}% smaller)"
    
    except Exception as e:
        return f"❌ Compression failed: {str(e)}"


@tool("images_to_pdf", "Combines multiple images into a single PDF file.")
def images_to_pdf(image_paths: str, output_path: str) -> str:
    """
    Combine multiple images into a PDF.
    
    Args:
        image_paths: Comma-separated list of image paths
        output_path: Path for the output PDF
        
    Returns:
        Path to the created PDF
    """
    if not PILLOW_AVAILABLE:
        return "❌ Pillow not installed"
    
    paths = [p.strip() for p in image_paths.split(',')]
    
    # Validate all paths
    for path in paths:
        if not os.path.exists(path):
            return f"❌ File not found: {path}"
    
    try:
        images = []
        for path in paths:
            img = Image.open(path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            images.append(img)
        
        if not images:
            return "❌ No valid images found"
        
        # Save as PDF
        first_image = images[0]
        if len(images) > 1:
            first_image.save(output_path, save_all=True, append_images=images[1:])
        else:
            first_image.save(output_path)
        
        return f"✓ Created PDF with {len(images)} pages: {output_path}"
    
    except Exception as e:
        return f"❌ PDF creation failed: {str(e)}"


# ============================================================================
# PDF TOOLS
# ============================================================================

@tool("pdf_to_text", "Extracts all text from a PDF file.")
def pdf_to_text(input_path: str, output_path: str = None) -> str:
    """
    Extract text from a PDF file.
    
    Args:
        input_path: Path to the PDF file
        output_path: Optional path to save text (if not provided, returns text)
        
    Returns:
        Extracted text or path to saved file
    """
    if not PYPDF2_AVAILABLE:
        return "❌ PyPDF2 not installed. Run: pip install PyPDF2"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    try:
        with open(input_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text_parts = []
            
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Page {page_num} ---\n{page_text}")
            
            full_text = "\n\n".join(text_parts)
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            return f"✓ Extracted text saved to: {output_path}"
        else:
            # Return truncated text
            if len(full_text) > 2000:
                return f"📄 PDF text ({len(reader.pages)} pages):\n{full_text[:2000]}...\n\n[Truncated - {len(full_text)} chars total]"
            return f"📄 PDF text ({len(reader.pages)} pages):\n{full_text}"
    
    except Exception as e:
        return f"❌ PDF extraction failed: {str(e)}"


@tool("merge_pdfs", "Merges multiple PDF files into one.")
def merge_pdfs(input_paths: str, output_path: str) -> str:
    """
    Merge multiple PDFs into a single file.
    
    Args:
        input_paths: Comma-separated list of PDF paths
        output_path: Path for the merged PDF
        
    Returns:
        Path to merged PDF
    """
    if not PYPDF2_AVAILABLE:
        return "❌ PyPDF2 not installed"
    
    paths = [p.strip() for p in input_paths.split(',')]
    
    for path in paths:
        if not os.path.exists(path):
            return f"❌ File not found: {path}"
    
    try:
        merger = PyPDF2.PdfMerger()
        
        for path in paths:
            merger.append(path)
        
        merger.write(output_path)
        merger.close()
        
        return f"✓ Merged {len(paths)} PDFs: {output_path}"
    
    except Exception as e:
        return f"❌ PDF merge failed: {str(e)}"


# ============================================================================
# DOCUMENT TOOLS
# ============================================================================

@tool("docx_to_text", "Extracts text from a Word document (.docx).")
def docx_to_text(input_path: str, output_path: str = None) -> str:
    """
    Extract text from a Word document.
    
    Args:
        input_path: Path to the .docx file
        output_path: Optional path to save text
        
    Returns:
        Extracted text or path to saved file
    """
    if not DOCX_AVAILABLE:
        return "❌ python-docx not installed. Run: pip install python-docx"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    try:
        doc = Document(input_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)
        
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(full_text)
            return f"✓ Extracted text saved to: {output_path}"
        else:
            if len(full_text) > 2000:
                return f"📄 Document text:\n{full_text[:2000]}...\n\n[Truncated - {len(full_text)} chars total]"
            return f"📄 Document text:\n{full_text}"
    
    except Exception as e:
        return f"❌ Document extraction failed: {str(e)}"


# ============================================================================
# PRESENTATION TOOLS
# ============================================================================

@tool("ppt_to_images", "Converts PowerPoint slides to image files (PNG).")
def ppt_to_images(input_path: str, output_folder: str = None) -> str:
    """
    Convert PowerPoint slides to images.
    
    Args:
        input_path: Path to the .pptx file
        output_folder: Folder to save images (default: same folder as pptx)
        
    Returns:
        Paths to created images
    """
    if not PPTX_AVAILABLE:
        return "❌ python-pptx not installed. Run: pip install python-pptx"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    # Note: python-pptx cannot directly export as images
    # This requires additional tools or workarounds
    
    # For Windows, we can use PowerPoint COM if available
    try:
        import subprocess
        
        # Try using PowerPoint via COM (Windows only)
        if sys.platform == 'win32':
            try:
                import win32com.client
                
                if not output_folder:
                    output_folder = os.path.dirname(input_path)
                
                os.makedirs(output_folder, exist_ok=True)
                
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                presentation = powerpoint.Presentations.Open(os.path.abspath(input_path), WithWindow=False)
                
                saved_files = []
                for i, slide in enumerate(presentation.Slides, 1):
                    output_path = os.path.join(output_folder, f"slide_{i:03d}.png")
                    slide.Export(output_path, "PNG")
                    saved_files.append(output_path)
                
                presentation.Close()
                powerpoint.Quit()
                
                return f"✓ Exported {len(saved_files)} slides to: {output_folder}"
            
            except ImportError:
                pass  # Fall through to alternative
        
        # Alternative: Extract embedded images and create text slides
        prs = Presentation(input_path)
        
        if not output_folder:
            output_folder = os.path.dirname(input_path)
        os.makedirs(output_folder, exist_ok=True)
        
        # Extract slide content as text files
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            
            output_path = os.path.join(output_folder, f"slide_{i:03d}_content.txt")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f"Slide {i}\n{'='*40}\n\n")
                f.write("\n\n".join(texts))
        
        return f"✓ Exported {len(prs.slides)} slide contents as text to: {output_folder}\n⚠️ Note: For image export, install pywin32 (pip install pywin32)"
    
    except Exception as e:
        return f"❌ PPT export failed: {str(e)}"


@tool("ppt_to_text", "Extracts all text from a PowerPoint presentation.")
def ppt_to_text(input_path: str) -> str:
    """
    Extract text from a PowerPoint presentation.
    
    Args:
        input_path: Path to the .pptx file
        
    Returns:
        All text from the presentation
    """
    if not PPTX_AVAILABLE:
        return "❌ python-pptx not installed"
    
    if not os.path.exists(input_path):
        return f"❌ File not found: {input_path}"
    
    try:
        prs = Presentation(input_path)
        all_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            all_text.append("\n".join(slide_texts))
        
        full_text = "\n\n".join(all_text)
        
        if len(full_text) > 2000:
            return f"📊 Presentation ({len(prs.slides)} slides):\n{full_text[:2000]}...\n\n[Truncated]"
        return f"📊 Presentation ({len(prs.slides)} slides):\n{full_text}"
    
    except Exception as e:
        return f"❌ PPT text extraction failed: {str(e)}"


# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def file_convert_status() -> dict:
    """Get status of file conversion capabilities."""
    return {
        "pillow": PILLOW_AVAILABLE,
        "pypdf2": PYPDF2_AVAILABLE,
        "python_pptx": PPTX_AVAILABLE,
        "python_docx": DOCX_AVAILABLE,
        "supported_image_formats": ["png", "jpg", "bmp", "gif", "webp", "ico", "tiff"],
    }
